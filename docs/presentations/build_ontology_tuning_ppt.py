from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "presentations"
ASSET_DIR = OUT_DIR / "ontology_tuning_assets"
PPTX_PATH = OUT_DIR / "PID智能整定_本体驱动闭环说明.pptx"
PREVIEW_DIR = OUT_DIR / "ontology_tuning_preview"


W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(7, 15, 26)
PANEL = RGBColor(15, 30, 46)
PANEL2 = RGBColor(20, 42, 64)
LINE = RGBColor(44, 71, 98)
TEXT = RGBColor(230, 242, 255)
MUTED = RGBColor(151, 178, 207)
BLUE = RGBColor(43, 133, 255)
CYAN = RGBColor(48, 214, 222)
GREEN = RGBColor(75, 245, 151)
YELLOW = RGBColor(246, 196, 79)


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.32, 8.6, 0.5, size=23, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.55, 0.86, 10.7, 0.32, size=10.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.28), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def panel(slide, x, y, w, h, fill=PANEL, line=LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    return sh


def chip(slide, text, x, y, w, color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(13, 45, 78)
    sh.line.color.rgb = color
    add_text(slide, text, x + 0.12, y + 0.07, w - 0.24, 0.16, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    return sh


def bullet_list(slide, items, x, y, w, line_h=0.42, size=14, color=TEXT):
    for i, item in enumerate(items):
        yy = y + i * line_h
        add_text(slide, "•", x, yy, 0.18, 0.25, size=size, color=CYAN, bold=True)
        add_text(slide, item, x + 0.25, yy, w - 0.25, 0.3, size=size, color=color)


def add_image_fit(slide, path, x, y, w, h):
    path = Path(path)
    img = Image.open(path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        height = h
        width = h * img_ratio
    else:
        width = w
        height = w / img_ratio
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_step(slide, num, title, desc, x, y, w=1.85):
    panel(slide, x, y, w, 1.0, fill=PANEL2)
    add_text(slide, str(num), x + 0.12, y + 0.18, 0.32, 0.25, size=17, color=GREEN, bold=True)
    add_text(slide, title, x + 0.48, y + 0.18, w - 0.55, 0.24, size=12.5, bold=True)
    add_text(slide, desc, x + 0.48, y + 0.54, w - 0.62, 0.3, size=8.6, color=MUTED)


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_text(s, "PID 智能整定如何结合本体", 0.65, 0.7, 7.2, 0.7, size=33, bold=True)
    add_text(s, "从历史数据画像、MCP 本体检索，到窗口策略、LLM 评审与整定闭环", 0.7, 1.55, 8.6, 0.35, size=15, color=MUTED)
    chip(s, "当前项目实现说明", 0.72, 2.1, 1.7, CYAN)
    panel(s, 0.7, 3.0, 3.2, 1.4)
    add_text(s, "核心目标", 0.95, 3.25, 1.6, 0.28, size=13, color=MUTED)
    add_text(s, "让窗口候选不只看曲线，而能理解工艺语义", 0.95, 3.68, 2.6, 0.42, size=18, bold=True)
    panel(s, 4.2, 3.0, 3.2, 1.4)
    add_text(s, "当前链路", 4.45, 3.25, 1.6, 0.28, size=13, color=MUTED)
    add_text(s, "LoopFeatures + MCP 本体 + LLM 策略 + 算法族", 4.45, 3.68, 2.6, 0.42, size=18, bold=True)
    add_image_fit(s, ASSET_DIR / "ppt_window_candidate_result.png", 7.7, 1.2, 5.0, 4.9)
    add_text(s, "截图来自当前 /monitoring 工业风驾驶舱", 7.75, 6.32, 3.8, 0.22, size=9, color=MUTED)

    # Slide 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "一条本体驱动的整定链路", "前端负责可解释流程展示；后端负责数据画像、本体检索、策略生成、算法族运行与事件沉淀。")
    add_step(s, 1, "数据画像", "LoopFeatures 计算原始特征", 0.62, 1.72)
    add_step(s, 2, "本体检索", "通过 MCP chat 查询回路知识", 2.68, 1.72)
    add_step(s, 3, "策略生成", "LLM 输出窗口算法策略 JSON", 4.74, 1.72)
    add_step(s, 4, "算法族运行", "MV 阶跃/斜坡/SP/稳态扰动", 6.8, 1.72)
    add_step(s, 5, "LLM 评审", "结合证据解释候选窗口", 8.86, 1.72)
    add_step(s, 6, "辨识整定", "多轮辨识、精修、PID 评估", 10.92, 1.72)
    panel(s, 0.7, 3.4, 5.8, 2.55)
    add_text(s, "后端关键模块", 0.98, 3.68, 2.4, 0.25, size=15, bold=True)
    bullet_list(s, [
        "core/shared/loop_features.py：历史数据画像与监控特征",
        "ontology_mcp_context.py：注册 MCP 工具检索本体上下文",
        "ontology_policy_builder.py：默认策略与回路类型先验",
        "window_policy_advisor.py：LLM 将本体/画像转为策略 JSON",
    ], 0.98, 4.15, 5.0, size=10.5, line_h=0.36)
    panel(s, 6.8, 3.4, 5.75, 2.55)
    add_text(s, "前端呈现重点", 7.08, 3.68, 2.4, 0.25, size=15, bold=True)
    bullet_list(s, [
        "窗口候选菜单按 6 步流程展示执行状态",
        "本体返回原文、策略字段、消费字段、LLM 评审均可追溯",
        "整定任务页展示准入校验、辨识、评审、整定、评估闭环",
        "MCP 服务配置页管理本体问答工具接入",
    ], 7.08, 4.15, 5.0, size=10.5, line_h=0.36)

    # Slide 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "MCP 服务配置：把本体工具注册进系统", "系统设置中配置服务名、Transport、服务地址和通用 MCP JSON；后端通过 /api/mcp-servers 读取。")
    add_image_fit(s, ASSET_DIR / "ppt_mcp_config.png", 0.7, 1.55, 7.3, 4.65)
    panel(s, 8.35, 1.55, 4.25, 4.65)
    add_text(s, "它在链路中承担什么？", 8.65, 1.9, 3.2, 0.25, size=16, bold=True)
    bullet_list(s, [
        "不是把本体写死在项目里，而是通过 MCP 工具查询。",
        "当前后端优先寻找启用服务中的 chat 工具。",
        "查询问题由回路位号和回路类型动态拼出。",
        "返回内容会作为 ontology_mcp_content_raw 注入窗口策略和评审。"
    ], 8.65, 2.45, 3.65, size=11.4, line_h=0.48)
    chip(s, "关键代码：core/mcp_config.py + api/mcp_config_routes.py", 8.62, 5.68, 3.65, CYAN)

    # Slide 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "本体查询问题：围绕窗口选择所需的工艺语义", "后端根据回路名与类型组织自然语言问题，要求本体返回变量角色、动态先验、扰动和窗口建议。")
    panel(s, 0.75, 1.55, 5.4, 4.9)
    add_text(s, "查询重点", 1.05, 1.86, 1.5, 0.25, size=16, bold=True)
    bullet_list(s, [
        "PV/CV、MV、SP、DV 的名称与物理含义",
        "对象阶次、正/反作用、过程增益方向",
        "时间尺度、滤波时间、死区、典型动态先验",
        "常见工况/扰动场景",
        "窗口辨识时应优先选择或避开的窗口特征",
    ], 1.05, 2.45, 4.8, size=13, line_h=0.48)
    panel(s, 6.55, 1.55, 5.95, 4.9)
    add_text(s, "为什么这一步重要", 6.85, 1.86, 2.2, 0.25, size=16, bold=True)
    bullet_list(s, [
        "仅靠 R² 或 NRMSE 容易把噪声/扰动片段当成好窗口。",
        "本体能告诉系统：这个回路的合理激励、时间尺度和增益方向。",
        "LLM 将本体事实转成可执行策略，再让算法族按策略筛选。",
        "最终 LLM 评审可以解释：为什么某窗口更合理，为什么某窗口有风险。",
    ], 6.85, 2.45, 5.1, size=13, line_h=0.48)

    # Slide 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "窗口候选页面：按流程显式展示智能体执行", "点击开始后，不直接吐窗口，而是按“画像 → 本体 → 策略 → 算法 → LLM → 准入”串行展示。")
    add_image_fit(s, ASSET_DIR / "ppt_window_candidate_start.png", 0.7, 1.45, 5.6, 4.85)
    add_image_fit(s, ASSET_DIR / "ppt_window_candidate_flow.png", 6.55, 1.45, 5.95, 4.85)
    add_text(s, "左：启动入口；右：流程运行中的步骤卡片和回路画像。", 0.78, 6.35, 6.5, 0.25, size=10, color=MUTED)

    # Slide 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "策略生成：LLM 不直接替代算法，而是给算法族下约束", "策略字段被分成“实际消费”和“展示/审计”，减少黑箱感。")
    panel(s, 0.7, 1.52, 4.05, 4.95)
    add_text(s, "可被算法消费的字段", 1.0, 1.86, 2.6, 0.25, size=16, bold=True)
    bullet_list(s, [
        "preferred/deprioritized/disabled_algorithm_families",
        "min_mv_excitation / min_sp_excitation",
        "max_mv_saturation_ratio / max_pv_noise_ratio",
        "min_window_duration_s / min_window_points",
        "pre_window_s / post_window_s",
        "steady_scan_window_s / steady_scan_step_s",
        "max_candidates_per_family",
    ], 1.0, 2.35, 3.45, size=10.6, line_h=0.36)
    panel(s, 4.95, 1.52, 3.6, 4.95)
    add_text(s, "主要用于解释/审计", 5.25, 1.86, 2.3, 0.25, size=16, bold=True)
    bullet_list(s, [
        "expected_gain_sign",
        "expected_time_constant_range_s",
        "expected_dead_time_range_s",
        "ontology_evidence",
        "rationale",
        "llm_policy_raw_text",
    ], 5.25, 2.35, 2.85, size=11.4, line_h=0.43)
    panel(s, 8.75, 1.52, 3.65, 4.95)
    add_text(s, "当前实现边界", 9.05, 1.86, 2.2, 0.25, size=16, bold=True)
    bullet_list(s, [
        "算法族仍是确定性 provider。",
        "LLM 策略影响算法族优先级和筛选门槛。",
        "LLM 最终评审解释候选窗口，但不应绕过准入规则。",
        "没有合格窗口时，应阻断正式辨识或降为诊断辨识。",
    ], 9.05, 2.35, 2.9, size=11.4, line_h=0.48)

    # Slide 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "LLM 评审：把本体证据和算法窗口合在一起解释", "评审不是只看分数，而是对照本体事实、LoopFeatures 和候选窗口逐项判断。")
    add_image_fit(s, ASSET_DIR / "ppt_window_candidate_llm.png", 0.7, 1.48, 7.6, 4.95)
    panel(s, 8.55, 1.48, 3.95, 4.95)
    add_text(s, "评审输出", 8.85, 1.82, 1.7, 0.25, size=16, bold=True)
    bullet_list(s, [
        "LLM 选中窗口与确定性窗口是否一致",
        "选择原因：激励强度、相关性、饱和、点数、工况",
        "引用的本体事实和来源",
        "候选窗口逐项评审：优先、可接受、风险",
        "是否允许进入正式辨识",
    ], 8.85, 2.32, 3.2, size=11.4, line_h=0.45)

    # Slide 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "闭环整定：本体驱动窗口只是第一关", "窗口选好后，系统仍会进入多轮辨识、模型评审、精修建议、PID 整定和性能评估。")
    panel(s, 0.75, 1.55, 11.75, 1.25)
    stages = ["候选窗口", "模型辨识", "LLM评审", "精修建议", "PID整定", "性能评估"]
    for idx, name in enumerate(stages):
        x = 1.0 + idx * 1.85
        chip(s, name, x, 2.0, 1.35, [CYAN, BLUE, GREEN, YELLOW, BLUE, GREEN][idx])
        if idx < len(stages) - 1:
            add_text(s, "→", x + 1.44, 2.02, 0.3, 0.2, size=15, color=MUTED, bold=True)
    panel(s, 0.75, 3.25, 5.55, 2.6)
    add_text(s, "后端闭环机制", 1.05, 3.56, 2.1, 0.25, size=16, bold=True)
    bullet_list(s, [
        "runner.py 支持多轮辨识 → LLM评审 → 精修建议 → 再辨识。",
        "identification_refinement_advisor.py 决定换窗口、缩模型池或给 L 初值。",
        "identification_advisor.py 只输出 accept/downgrade，不直接 reject 中止。",
    ], 1.05, 4.05, 4.85, size=11.3, line_h=0.45)
    panel(s, 6.6, 3.25, 5.9, 2.6)
    add_text(s, "产品化建议", 6.9, 3.56, 2.1, 0.25, size=16, bold=True)
    bullet_list(s, [
        "保留“诊断辨识”和“正式辨识”两种准入结论。",
        "把本体返回原文、策略 JSON、消费字段、LLM 原文全部沉淀到事件。",
        "对没有合格窗口的回路，优先输出数据采集/试验建议，而不是强行整定。",
    ], 6.9, 4.05, 5.0, size=11.3, line_h=0.45)

    prs.save(PPTX_PATH)


def render_preview_images():
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    font_path = "C:/Windows/Fonts/msyh.ttc"
    try:
        title_font = ImageFont.truetype(font_path, 38)
        body_font = ImageFont.truetype(font_path, 22)
        small_font = ImageFont.truetype(font_path, 16)
    except Exception:
        title_font = body_font = small_font = ImageFont.load_default()
    slide_titles = [
        "PID 智能整定如何结合本体",
        "一条本体驱动的整定链路",
        "MCP 服务配置：把本体工具注册进系统",
        "本体查询问题：围绕窗口选择所需的工艺语义",
        "窗口候选页面：按流程显式展示智能体执行",
        "策略生成：LLM 不直接替代算法",
        "LLM 评审：结合本体证据解释窗口",
        "闭环整定：本体驱动窗口只是第一关",
    ]
    for i, title in enumerate(slide_titles, start=1):
        img = Image.new("RGB", (1280, 720), (7, 15, 26))
        d = ImageDraw.Draw(img)
        d.text((48, 36), title, fill=(230, 242, 255), font=title_font)
        d.line((48, 96, 1230, 96), fill=(44, 71, 98), width=2)
        if i == 1:
            d.text((60, 160), "可编辑 PPT 已生成；预览图用于快速核对页序和主题。", fill=(151, 178, 207), font=body_font)
        elif i in {3, 5, 7}:
            screenshot = {
                3: ASSET_DIR / "ppt_mcp_config.png",
                5: ASSET_DIR / "ppt_window_candidate_flow.png",
                7: ASSET_DIR / "ppt_window_candidate_llm.png",
            }[i]
            shot = Image.open(screenshot).convert("RGB")
            shot.thumbnail((980, 520))
            img.paste(shot, (150, 135))
        else:
            d.text((60, 160), "本页为文字、流程图和可编辑形状说明页。", fill=(151, 178, 207), font=body_font)
            d.text((60, 220), "请打开 PPTX 查看完整排版。", fill=(75, 245, 151), font=small_font)
        img.save(PREVIEW_DIR / f"slide_{i:02d}.png")


if __name__ == "__main__":
    build_deck()
    render_preview_images()
    print(PPTX_PATH)
    print(PREVIEW_DIR)
